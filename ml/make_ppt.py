"""
make_ppt.py — generates HealthEase_Presentation.pptx (24 slides, premium-simple).
Run: venv/bin/python make_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
CLEAN = os.path.join(HERE, "figures_clean")
SHOTS = os.path.join(HERE, "screenshots")
OUT = os.path.join(HERE, "..", "HealthEase_Presentation.pptx")

# ---- palette ----
TEAL = RGBColor(0x0D, 0x94, 0x88)
BLUE = RGBColor(0x25, 0x63, 0xEB)
HI = RGBColor(0x14, 0xB8, 0xA6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD = RGBColor(0xF1, 0xF5, 0xF9)

FONT = "Poppins"
BODYFONT = "Inter"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def solid(shape, color, line=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 1)
    shape.shadow.inherit = False


def no_autofit(tf):
    tf.word_wrap = True


def set_text(tf, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, font)."""
    no_autofit(tf)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        for (text, size, color, bold, font) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font


def textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = slide.shapes.add_textbox(l, t, w, h)
    set_text(tb.text_frame, runs, align, anchor, space)
    return tb


def box(slide, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    solid(sp, fill, line, line_w)
    return sp


def boxtext(slide, l, t, w, h, fill, runs, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            line=None, line_w=None, space=3):
    sp = box(slide, l, t, w, h, fill, shape, line, line_w)
    tf = sp.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    set_text(tf, runs, align, anchor, space)
    return sp


def arrow(slide, x1, y1, x2, y2, color=BLUE, w=2.0):
    cn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    line_elem = cn.line._get_or_add_ln()
    tail = line_elem.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line_elem.append(tail)
    cn.shadow.inherit = False
    return cn


def base(section=None, divider=False):
    """new slide with bg + top accent + footer."""
    s = prs.slides.add_slide(BLANK)
    bg = box(s, 0, 0, SW, SH, BG, MSO_SHAPE.RECTANGLE)
    # footer bar
    box(s, 0, SH - Inches(0.32), SW, Inches(0.32), TEAL, MSO_SHAPE.RECTANGLE)
    textbox(s, Inches(0.4), SH - Inches(0.34), Inches(4), Inches(0.3),
            [[("HealthEase — AI Healthcare Platform", 9, WHITE, False, BODYFONT)]],
            anchor=MSO_ANCHOR.MIDDLE)
    n = len(prs.slides._sldIdLst)
    textbox(s, SW - Inches(1.2), SH - Inches(0.34), Inches(0.8), Inches(0.3),
            [[(str(n), 9, WHITE, True, BODYFONT)]], align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE)
    return s


def header(s, title, kicker=None):
    box(s, Inches(0.4), Inches(0.45), Inches(0.14), Inches(0.55), TEAL, MSO_SHAPE.RECTANGLE)
    runs = []
    if kicker:
        runs.append([(kicker.upper(), 11, HI, True, BODYFONT)])
    runs.append([(title, 28, SLATE, True, FONT)])
    textbox(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.95), runs, space=2)


def bullets(items, size=15, color=SLATE, marker="•  ", mcolor=TEAL):
    out = []
    for it in items:
        out.append([(marker, size, mcolor, True, BODYFONT), (it, size, color, False, BODYFONT)])
    return out


# =========================================================== SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, SW, SH, SLATE, MSO_SHAPE.RECTANGLE)
box(s, 0, 0, SW, Inches(0.25), TEAL, MSO_SHAPE.RECTANGLE)
box(s, 0, SH - Inches(0.25), SW, Inches(0.25), TEAL, MSO_SHAPE.RECTANGLE)
boxtext(s, SW/2 - Inches(0.6), Inches(1.1), Inches(1.2), Inches(1.2), TEAL,
        [[("✚", 44, WHITE, True, FONT)]], shape=MSO_SHAPE.OVAL)
textbox(s, Inches(1), Inches(2.6), Inches(11.33), Inches(1.2),
        [[("HealthEase", 60, WHITE, True, FONT)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(3.8), Inches(11.33), Inches(0.7),
        [[("AI-Powered Healthcare Management Platform", 24, HI, True, FONT)]],
        align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.5), Inches(11.33), Inches(0.5),
        [[("Final Year B.Tech Project", 16, RGBColor(0xCB,0xD5,0xE1), False, BODYFONT)]],
        align=PP_ALIGN.CENTER)
boxtext(s, SW - Inches(5.0), Inches(5.4), Inches(4.6), Inches(1.7),
        RGBColor(0x33,0x41,0x55),
        [[("Presented by", 11, RGBColor(0x94,0xA3,0xB8), True, BODYFONT)],
         [("2022BITE018  —  Vardan Pal", 12.5, WHITE, True, BODYFONT)],
         [("2022BITE037  —  U Rajeshwar", 12.5, WHITE, True, BODYFONT)],
         [("2022BITE007  —  Karan Kharadi", 12.5, WHITE, True, BODYFONT)]],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, space=6)

# =========================================================== SLIDE 2 — AGENDA
s = base(); header(s, "Agenda", "Outline")
secs = ["Introduction", "Background & Literature", "System Design",
        "Implementation", "Results & Testing", "Conclusion"]
for i, name in enumerate(secs):
    col = i // 3; row = i % 3
    x = Inches(0.9 + col * 6.2); y = Inches(1.7 + row * 1.5)
    boxtext(s, x, y, Inches(0.8), Inches(0.8), TEAL,
            [[(str(i+1), 24, WHITE, True, FONT)]], shape=MSO_SHAPE.OVAL)
    boxtext(s, x + Inches(1.0), y, Inches(4.6), Inches(0.8), CARD,
            [[(name, 17, SLATE, True, BODYFONT)]], align=PP_ALIGN.LEFT)

# =========================================================== SLIDE 3 — ABSTRACT
s = base(); header(s, "Abstract", "Overview")
boxtext(s, Inches(1.1), Inches(1.9), Inches(11.1), Inches(3.6), CARD,
        [[("HealthEase is a full-stack MERN platform that digitises paper and "
           "handwritten prescriptions using an AI vision-language OCR engine, then "
           "applies a custom-trained machine-learning classifier to map each "
           "prescribed medicine to the disease it treats.", 17, SLATE, False, BODYFONT)],
         [("", 6, SLATE, False, BODYFONT)],
         [("It tracks patient vitals and medication compliance, computes a Smart "
           "Health Score (0–100), checks drug interactions, and connects patients to "
           "doctors through telemedicine — unifying fragmented healthcare workflows "
           "into one secure system.", 17, SLATE, False, BODYFONT)]],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, space=8)

# =========================================================== SLIDE 4 — PROBLEM
s = base(); header(s, "Problem Statement", "Why HealthEase")
probs = [
    ("📄", "Hard-to-read prescriptions", "Paper and handwritten prescriptions are difficult to digitise."),
    ("❓", "Unclear medicines", "Patients don't understand what each medicine is for."),
    ("💊", "Poor compliance", "Missed doses and refills lead to worse outcomes."),
    ("🧩", "Fragmented data", "Vitals, prescriptions and consultations live in silos."),
]
for i, (ic, t, d) in enumerate(probs):
    col = i % 2; row = i // 2
    x = Inches(0.9 + col * 6.0); y = Inches(1.8 + row * 2.3)
    c = box(s, x, y, Inches(5.5), Inches(2.0), CARD)
    boxtext(s, x + Inches(0.25), y + Inches(0.3), Inches(0.9), Inches(0.9), AMBER,
            [[(ic, 22, WHITE, True, FONT)]], shape=MSO_SHAPE.OVAL)
    textbox(s, x + Inches(1.35), y + Inches(0.3), Inches(3.9), Inches(1.5),
            [[(t, 16, SLATE, True, BODYFONT)], [(d, 12, GRAY, False, BODYFONT)]], space=4)

# =========================================================== SLIDE 5 — MOTIVATION
s = base(); header(s, "Motivation", "Background")
textbox(s, Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.5),
        bullets([
            "Rising chronic-disease burden demands better tracking.",
            "Prescriptions must be made understandable to patients.",
            "Compliance monitoring should be active, not passive.",
            "Healthcare data needs a single, unified platform.",
            "Modern AI (vision models + ML) enables real clinical value.",
        ], size=16), space=12)
boxtext(s, Inches(8.4), Inches(1.9), Inches(4.2), Inches(3.8), TEAL,
        [[("🏥", 60, WHITE, True, FONT)],
         [("Patient-centric, AI-driven care", 14, WHITE, True, BODYFONT)]],
        anchor=MSO_ANCHOR.MIDDLE, space=10)

# =========================================================== SLIDE 6 — OBJECTIVES
s = base(); header(s, "Objectives", "Goals")
objs = [
    "Automatically extract data from prescription images (OCR).",
    "Map medicines to their indication using a trained ML model.",
    "Track vitals and medication compliance with analytics.",
    "Quantify patient health as a single Smart Health Score.",
    "Enable secure multi-role telemedicine (Patient/Doctor/Admin).",
    "Deliver everything in one responsive, secure web app.",
]
for i, o in enumerate(objs):
    y = Inches(1.65 + i * 0.85)
    boxtext(s, Inches(0.9), y, Inches(0.62), Inches(0.62), BLUE if i % 2 else TEAL,
            [[(str(i+1), 18, WHITE, True, FONT)]], shape=MSO_SHAPE.OVAL)
    boxtext(s, Inches(1.7), y, Inches(10.7), Inches(0.62), CARD,
            [[(o, 15, SLATE, False, BODYFONT)]], align=PP_ALIGN.LEFT)

# =========================================================== SLIDE 7 — LITERATURE
s = base(); header(s, "Literature Survey", "Existing Systems")
rows = [
    ["Feature", "HealthEase", "Practo", "1mg", "EHR"],
    ["OCR prescription reading", "✓", "✗", "✗", "✗"],
    ["AI medicine → disease mapping", "✓", "✗", "~", "✗"],
    ["Smart Health Score", "✓", "✗", "✗", "✗"],
    ["Compliance tracking", "✓", "✗", "✓", "~"],
    ["Vitals analytics", "✓", "~", "✗", "✓"],
    ["Drug-interaction check", "✓", "✗", "✓", "~"],
    ["Real-time notifications", "✓", "✓", "✓", "✗"],
]
t = s.shapes.add_table(len(rows), 5, Inches(0.7), Inches(1.65),
                       Inches(11.9), Inches(5.2)).table
t.columns[0].width = Inches(4.3)
for c in range(1, 5):
    t.columns[c].width = Inches(1.9)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = t.cell(ri, ci)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.name = BODYFONT
        if ri == 0:
            r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = TEAL if ci == 1 else SLATE
        else:
            r.font.size = Pt(12.5)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE2,0xF5,0xF2) if ci == 1 else (WHITE if ri % 2 else CARD)
            if val == "✓":
                r.font.color.rgb = GREEN; r.font.bold = True
            elif val == "✗":
                r.font.color.rgb = RGBColor(0xCB,0xD5,0xE1)
            elif val == "~":
                r.font.color.rgb = AMBER; r.font.bold = True
            else:
                r.font.color.rgb = SLATE
                if ci == 0: r.font.bold = True

# =========================================================== SLIDE 8 — TECH STACK
s = base(); header(s, "Technology Stack", "System Design")
bands = [
    ("Frontend", "React + Vite · Tailwind CSS · Context API", TEAL),
    ("Backend", "Node.js + Express · JWT · bcrypt", RGBColor(0x0E,0x7A,0x9B)),
    ("Database", "MongoDB + Mongoose", RGBColor(0x15,0x65,0xC0)),
    ("Real-time", "Socket.IO live notifications", BLUE),
    ("OCR Microservice", "Python + FastAPI", RGBColor(0x0D,0x94,0x88)),
    ("OCR Engine", "Groq Llama-4 Vision (pre-trained)", HI),
    ("Machine Learning", "scikit-learn (custom-trained classifier)", GREEN),
    ("PDF Export", "jsPDF + html2canvas", RGBColor(0x64,0x74,0x8B)),
]
y0 = 1.62; bh = 0.62; gap = 0.05
for i, (name, desc, col) in enumerate(bands):
    y = Inches(y0 + i * (bh + gap))
    boxtext(s, Inches(0.9), y, Inches(3.4), Inches(bh), col,
            [[(name, 14, WHITE, True, BODYFONT)]], align=PP_ALIGN.LEFT)
    boxtext(s, Inches(4.4), y, Inches(8.0), Inches(bh), CARD,
            [[(desc, 13, SLATE, False, BODYFONT)]], align=PP_ALIGN.LEFT)

# =========================================================== SLIDE 9 — ARCHITECTURE
s = base(); header(s, "System Architecture", "System Design")
client = boxtext(s, Inches(4.4), Inches(1.55), Inches(4.5), Inches(0.85), TEAL,
        [[("React + Vite Client", 15, WHITE, True, BODYFONT)]])
gw = boxtext(s, Inches(4.4), Inches(3.0), Inches(4.5), Inches(0.85), BLUE,
        [[("Node.js + Express API Gateway", 14, WHITE, True, BODYFONT)]])
db = boxtext(s, Inches(1.5), Inches(4.6), Inches(4.0), Inches(0.85), SLATE,
        [[("MongoDB Database", 14, WHITE, True, BODYFONT)]])
py = boxtext(s, Inches(7.9), Inches(4.6), Inches(4.0), Inches(0.85), GREEN,
        [[("Python FastAPI Service", 14, WHITE, True, BODYFONT)]])
ocr = boxtext(s, Inches(7.0), Inches(6.0), Inches(2.7), Inches(0.7), HI,
        [[("Groq Vision OCR", 12, WHITE, True, BODYFONT)]])
mlb = boxtext(s, Inches(10.0), Inches(6.0), Inches(2.7), Inches(0.7), RGBColor(0x16,0xA3,0x4A),
        [[("ML Classifier", 12, WHITE, True, BODYFONT)]])
arrow(s, Inches(6.65), Inches(2.4), Inches(6.65), Inches(3.0), BLUE)
arrow(s, Inches(3.5), Inches(3.85), Inches(3.5), Inches(4.6), TEAL)
arrow(s, Inches(9.9), Inches(3.85), Inches(9.9), Inches(4.6), GREEN)
arrow(s, Inches(8.35), Inches(5.45), Inches(8.35), Inches(6.0), HI)
arrow(s, Inches(11.3), Inches(5.45), Inches(11.3), Inches(6.0), GREEN)
textbox(s, Inches(2.7), Inches(3.9), Inches(2.5), Inches(0.4),
        [[("Mongoose ODM", 10, GRAY, True, BODYFONT)]])
textbox(s, Inches(9.95), Inches(3.9), Inches(2.5), Inches(0.4),
        [[("HTTP REST", 10, GRAY, True, BODYFONT)]])
textbox(s, Inches(6.8), Inches(2.5), Inches(3.2), Inches(0.4),
        [[("REST + WebSockets", 10, GRAY, True, BODYFONT)]])

# =========================================================== SLIDE 10 — DFD
s = base(); header(s, "Data Flow Diagram (Level 1)", "System Design")
# entities
boxtext(s, Inches(0.6), Inches(2.0), Inches(1.9), Inches(0.8), SLATE,
        [[("Patient", 13, WHITE, True, BODYFONT)]], shape=MSO_SHAPE.RECTANGLE)
boxtext(s, Inches(0.6), Inches(5.0), Inches(1.9), Inches(0.8), SLATE,
        [[("Doctor", 13, WHITE, True, BODYFONT)]], shape=MSO_SHAPE.RECTANGLE)
procs = [
    ("1.0 Upload\nPrescription", 3.1, 1.7, TEAL),
    ("2.0 OCR\nExtraction", 5.6, 1.7, TEAL),
    ("3.0 Medicine\nClassification", 8.1, 1.7, BLUE),
    ("4.0 Vitals\nLogging", 3.1, 4.0, TEAL),
    ("5.0 Health\nScore", 5.6, 4.0, GREEN),
    ("6.0 Tele-\nmedicine", 8.1, 4.0, BLUE),
]
for txt, x, y, c in procs:
    boxtext(s, Inches(x), Inches(y), Inches(2.0), Inches(1.2), c,
            [[(txt, 12, WHITE, True, BODYFONT)]], shape=MSO_SHAPE.OVAL)
# data stores
for txt, x, y in [("Prescriptions", 10.8, 1.9), ("Medicines", 10.8, 3.1),
                  ("Vitals", 10.8, 4.3)]:
    boxtext(s, Inches(x), Inches(y), Inches(2.2), Inches(0.6), CARD,
            [[(txt, 12, SLATE, True, BODYFONT)]], shape=MSO_SHAPE.RECTANGLE,
            line=TEAL, line_w=1.5)
arrow(s, Inches(2.5), Inches(2.3), Inches(3.1), Inches(2.3), TEAL)
arrow(s, Inches(5.1), Inches(2.3), Inches(5.6), Inches(2.3), TEAL)
arrow(s, Inches(7.6), Inches(2.3), Inches(8.1), Inches(2.3), BLUE)
arrow(s, Inches(2.5), Inches(5.0), Inches(3.1), Inches(4.6), TEAL)
arrow(s, Inches(5.1), Inches(4.6), Inches(5.6), Inches(4.6), GREEN)
arrow(s, Inches(7.6), Inches(4.6), Inches(8.1), Inches(4.6), BLUE)
arrow(s, Inches(10.1), Inches(2.3), Inches(10.8), Inches(2.2), GRAY, 1.5)
arrow(s, Inches(10.1), Inches(2.3), Inches(10.8), Inches(3.4), GRAY, 1.5)

# =========================================================== SLIDE 11 — USE CASE
s = base(); header(s, "Use-Case Diagram", "System Design")
box(s, Inches(3.4), Inches(1.55), Inches(6.5), Inches(5.2), WHITE,
    MSO_SHAPE.RECTANGLE, line=TEAL, line_w=1.5)
textbox(s, Inches(3.4), Inches(1.6), Inches(6.5), Inches(0.4),
        [[("HealthEase System", 13, TEAL, True, BODYFONT)]], align=PP_ALIGN.CENTER)
# actors
for txt, x, y in [("Patient", 0.7, 3.3), ("Doctor", 11.4, 3.3), ("Admin", 6.0, 6.85)]:
    boxtext(s, Inches(x), Inches(y), Inches(1.3), Inches(0.6), SLATE,
            [[("🧍 " + txt, 12, WHITE, True, BODYFONT)]], shape=MSO_SHAPE.RECTANGLE)
uc = [("Upload Prescription", 3.7, 2.1), ("Log Vitals", 3.7, 2.85),
      ("View Health Score", 3.7, 3.6), ("Book Consultation", 3.7, 4.35),
      ("Set Reminders", 3.7, 5.1),
      ("View Records", 7.3, 2.1), ("Add Diagnosis", 7.3, 2.85),
      ("Conduct Consult", 7.3, 3.6), ("Approve Doctors", 5.5, 5.5),
      ("Manage Users", 7.3, 5.5)]
for txt, x, y in uc:
    boxtext(s, Inches(x), Inches(y), Inches(2.5), Inches(0.6), TEAL,
            [[(txt, 11, WHITE, True, BODYFONT)]], shape=MSO_SHAPE.OVAL)

# =========================================================== SLIDE 12 — ER
s = base(); header(s, "Database Design (ER Diagram)", "System Design")
ent = [
    ("User", "email · role · passwordHash", 5.4, 1.55, TEAL),
    ("Patient", "age · gender · contact", 1.2, 3.0, BLUE),
    ("Doctor", "specialty · approved", 9.6, 3.0, BLUE),
    ("Prescription", "doctorName · medications[]", 1.2, 5.0, GREEN),
    ("Medicine", "dosage · frequency · stock", 5.4, 5.0, GREEN),
    ("Consultation", "date · status · notes", 9.6, 5.0, GREEN),
]
for name, fields, x, y, c in ent:
    boxtext(s, Inches(x), Inches(y), Inches(2.8), Inches(0.5), c,
            [[(name, 13, WHITE, True, BODYFONT)]], shape=MSO_SHAPE.RECTANGLE)
    boxtext(s, Inches(x), Inches(y + 0.5), Inches(2.8), Inches(0.7), CARD,
            [[(fields, 10, SLATE, False, BODYFONT)]], shape=MSO_SHAPE.RECTANGLE,
            line=c, line_w=1)
arrow(s, Inches(6.0), Inches(2.25), Inches(2.6), Inches(3.0), GRAY, 1.5)
arrow(s, Inches(7.2), Inches(2.25), Inches(10.6), Inches(3.0), GRAY, 1.5)
arrow(s, Inches(2.6), Inches(3.7), Inches(2.6), Inches(5.0), GRAY, 1.5)
arrow(s, Inches(4.0), Inches(5.35), Inches(5.4), Inches(5.35), GRAY, 1.5)
arrow(s, Inches(10.6), Inches(3.7), Inches(10.6), Inches(5.0), GRAY, 1.5)
textbox(s, Inches(0.6), Inches(2.5), Inches(2), Inches(0.3), [[("1 : 1", 9, GRAY, True, BODYFONT)]])
textbox(s, Inches(2.7), Inches(4.3), Inches(2), Inches(0.3), [[("1 : N", 9, GRAY, True, BODYFONT)]])

# =========================================================== SLIDE 13 — SEQUENCE
s = base(); header(s, "Sequence Diagram — Prescription Flow", "System Design")
lifelines = ["Patient", "Client", "API", "OCR Svc", "ML Model", "Database"]
xs = [1.2, 3.0, 4.9, 6.9, 8.9, 11.0]
for name, x in zip(lifelines, xs):
    boxtext(s, Inches(x - 0.6), Inches(1.55), Inches(1.7), Inches(0.55),
            TEAL, [[(name, 11, WHITE, True, BODYFONT)]])
    ln = s.shapes.add_connector(1, Inches(x + 0.25), Inches(2.1), Inches(x + 0.25), Inches(6.7))
    ln.line.color.rgb = RGBColor(0xCB,0xD5,0xE1); ln.line.width = Pt(1)
    ln.shadow.inherit = False
msgs = [
    (0, 1, "upload image"), (1, 2, "POST /api/ocr"), (2, 3, "forward image"),
    (3, 3, "vision request → Groq"), (3, 2, "extracted text"),
    (2, 4, "POST /classify-medicines"), (4, 2, "class + indication"),
    (2, 5, "save prescription"), (2, 1, "result"), (1, 0, "display"),
]
y = 2.45
for a, b, txt in msgs:
    xa, xb = xs[a] + 0.25, xs[b] + 0.25
    if a == b:
        textbox(s, Inches(xa + 0.05), Inches(y - 0.18), Inches(3), Inches(0.3),
                [[("↻ " + txt, 9.5, BLUE, True, BODYFONT)]])
    else:
        arrow(s, Inches(xa), Inches(y), Inches(xb), Inches(y),
              BLUE if xb > xa else GREEN, 1.5)
        mid = min(xa, xb)
        textbox(s, Inches(mid + 0.05), Inches(y - 0.28), Inches(abs(xb-xa)+1.5), Inches(0.3),
                [[(txt, 9.5, SLATE, False, BODYFONT)]])
    y += 0.42

# =========================================================== SLIDE 14 — OCR
s = base(); header(s, "OCR Prescription Reader", "Implementation")
textbox(s, Inches(0.8), Inches(1.75), Inches(5.4), Inches(4.5),
        bullets([
            "Type: AI Vision-Language OCR (not classic Tesseract).",
            "Engine: Groq Llama-4 Scout vision model (pre-trained).",
            "Handles printed AND handwritten prescriptions.",
            "Returns structured fields ready for the database.",
        ], size=15), space=12)
steps = ["Image Upload", "Preprocess\n(resize + RGB)", "Base64 Encode",
         "Vision Model\n+ Medical Prompt", "Structured Output\nDoctor/Diagnosis/Meds"]
y = 1.8
for i, st in enumerate(steps):
    boxtext(s, Inches(7.0), Inches(y), Inches(5.2), Inches(0.72),
            TEAL if i % 2 == 0 else BLUE, [[(st, 12.5, WHITE, True, BODYFONT)]])
    if i < len(steps) - 1:
        arrow(s, Inches(9.6), Inches(y + 0.72), Inches(9.6), Inches(y + 0.95), BLUE)
    y += 0.97

# =========================================================== SLIDE 15 — ML OVERVIEW
s = base(); header(s, "Medicine-Indication Classifier", "Implementation · ML")
textbox(s, Inches(0.8), Inches(1.75), Inches(5.6), Inches(4.5),
        bullets([
            "Predicts therapeutic class + disease from the medicine NAME.",
            "Generalises to drug names never seen in training.",
            "Built on WHO INN naming stems (e.g. -pril, -statin).",
            "Custom-trained — fully reproducible from our dataset.",
        ], size=15), space=12)
stem = [["Stem", "Class", "Example"],
        ["-pril", "ACE Inhibitor", "Lisinopril"],
        ["-statin", "Statin", "Atorvastatin"],
        ["-cillin", "Penicillin", "Amoxicillin"],
        ["-floxacin", "Fluoroquinolone", "Ciprofloxacin"],
        ["-prazole", "PPI", "Omeprazole"],
        ["-dipine", "Ca-Channel Blocker", "Amlodipine"]]
t = s.shapes.add_table(len(stem), 3, Inches(6.8), Inches(1.75),
                       Inches(5.7), Inches(4.4)).table
for ri, row in enumerate(stem):
    for ci, val in enumerate(row):
        cell = t.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
        r.font.name = BODYFONT; r.font.size = Pt(12.5)
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
        else:
            r.font.color.rgb = SLATE
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else CARD
            if ci == 0: r.font.bold = True; r.font.color.rgb = BLUE

# =========================================================== SLIDE 16 — ML METHOD
s = base(); header(s, "ML Methodology", "Implementation · ML")
stages = [
    "Curated Dataset — 321 medicines · 27 classes",
    "Feature Extraction — Character n-gram TF-IDF (n = 2–4)",
    "Train / Test Split — 75% / 25% stratified",
    "Model Training — Logistic Regression",
    "5-Fold Cross-Validation",
    "Evaluation — Accuracy · F1 · Confusion Matrix",
    "Deployed Model — medicine_classifier.joblib",
]
y = 1.6
for i, st in enumerate(stages):
    boxtext(s, Inches(3.0), Inches(y), Inches(7.3), Inches(0.62),
            TEAL if i % 2 == 0 else BLUE, [[(st, 13, WHITE, True, BODYFONT)]])
    if i < len(stages) - 1:
        arrow(s, Inches(6.65), Inches(y + 0.62), Inches(6.65), Inches(y + 0.74), BLUE, 2)
    y += 0.74

# =========================================================== SLIDE 17 — ML RESULTS
s = base(); header(s, "ML Results", "Implementation · ML")
cards = [("85.4%", "Cross-Validated Accuracy"), ("76.5%", "Test Accuracy"),
         ("0.78", "Macro F1-Score")]
for i, (big, lab) in enumerate(cards):
    x = Inches(0.8 + i * 2.15)
    c = box(s, x, Inches(1.7), Inches(1.95), Inches(1.5), CARD)
    textbox(s, x, Inches(1.8), Inches(1.95), Inches(0.8),
            [[(big, 30, GREEN, True, FONT)]], align=PP_ALIGN.CENTER)
    textbox(s, x, Inches(2.65), Inches(1.95), Inches(0.5),
            [[(lab, 10.5, SLATE, True, BODYFONT)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(0.8), Inches(3.5), Inches(6.2), Inches(0.4),
        [[("Model comparison (5-fold CV):", 13, SLATE, True, BODYFONT)]])
comp = [["Logistic Regression", "82.9%"], ["Linear SVM", "81.7%"], ["Naive Bayes", "77.9%"]]
for i, (m, v) in enumerate(comp):
    y = Inches(4.0 + i * 0.62)
    boxtext(s, Inches(0.8), y, Inches(4.2), Inches(0.52), CARD,
            [[(m, 12.5, SLATE, True, BODYFONT)]], align=PP_ALIGN.LEFT)
    boxtext(s, Inches(5.1), y, Inches(1.9), Inches(0.52), TEAL,
            [[(v, 12.5, WHITE, True, BODYFONT)]])
textbox(s, Inches(0.8), Inches(6.1), Inches(6.3), Inches(0.8),
        [[("The model independently rediscovered medical naming conventions from data.",
           12, BLUE, True, BODYFONT)]])
img = os.path.join(CLEAN, "g1_learning_curve.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(7.2), Inches(2.0), width=Inches(5.7))

# =========================================================== SLIDE 18 — FEATURES
s = base(); header(s, "Key Features", "Implementation")
feats = [
    ("🔐", "Multi-role Auth"), ("📄", "AI OCR Digitiser"), ("🧠", "Medicine→Disease"),
    ("⚠️", "Interaction Check"), ("📈", "Vitals Analytics"), ("💯", "Health Score"),
    ("💊", "Tracker+Reminders"), ("🔔", "Notifications"), ("🏥", "Telemedicine"),
    ("🛡️", "Admin Dashboard"), ("📑", "PDF Export"), ("🌗", "Dark Mode / UI"),
]
for i, (ic, t) in enumerate(feats):
    col = i % 4; row = i // 4
    x = Inches(0.7 + col * 3.1); y = Inches(1.75 + row * 1.55)
    box(s, x, y, Inches(2.85), Inches(1.35), CARD)
    boxtext(s, x + Inches(0.2), y + Inches(0.25), Inches(0.85), Inches(0.85),
            TEAL, [[(ic, 20, WHITE, True, FONT)]], shape=MSO_SHAPE.OVAL)
    textbox(s, x + Inches(1.15), y + Inches(0.45), Inches(1.6), Inches(0.7),
            [[(t, 12.5, SLATE, True, BODYFONT)]])

# =========================================================== SLIDE 19 — WORKFLOW
s = base(); header(s, "End-to-End Workflow", "Implementation")
flow = ["Upload\nPrescription", "OCR\nExtraction", "ML\nClassification",
        "Interaction\nCheck", "Schedule\nReminders", "Log\nVitals",
        "Update Health\nScore", "Doctor Review\n/ PDF Export"]
for i, st in enumerate(flow):
    col = i % 4; row = i // 4
    x = Inches(0.7 + col * 3.15); y = Inches(2.0 + row * 2.2)
    boxtext(s, x, y, Inches(2.5), Inches(1.4),
            TEAL if i % 2 == 0 else BLUE,
            [[(str(i+1), 18, WHITE, True, FONT)], [(st, 12, WHITE, True, BODYFONT)]],
            shape=MSO_SHAPE.OVAL, space=2)
    if col < 3:
        arrow(s, x + Inches(2.5), y + Inches(0.7), x + Inches(3.15), y + Inches(0.7), GRAY, 2)
    elif row == 0:
        arrow(s, x + Inches(1.25), y + Inches(1.4), x + Inches(1.25), y + Inches(2.2), GRAY, 2)

# =========================================================== SLIDE 20 — APP RESULTS
s = base(); header(s, "Application Results", "Results & Testing")
shots = [("Patient Dashboard", "dashboard.png"),
         ("Prescription Records", "prescriptions.png"),
         ("Prescription Digitizer", "upload.png"),
         ("Doctor Marketplace", "doctors.png"),
         ("Medicine Tracker", "medicine_tracker.png"),
         ("Vitals Analytics", "vitals.png")]
for i, (lab, fn) in enumerate(shots):
    col = i % 3; row = i // 3
    x = Inches(0.55 + col * 4.25); y = Inches(1.65 + row * 2.55)
    p = os.path.join(SHOTS, fn)
    box(s, x - Inches(0.04), y - Inches(0.04), Inches(4.0), Inches(2.32),
        WHITE, MSO_SHAPE.RECTANGLE, line=TEAL, line_w=1.25)
    if os.path.exists(p):
        s.shapes.add_picture(p, x, y, width=Inches(3.92))
    textbox(s, x, y + Inches(2.0), Inches(3.92), Inches(0.4),
            [[(lab, 11, SLATE, True, BODYFONT)]], align=PP_ALIGN.CENTER)

# =========================================================== SLIDE 21 — ML GRAPHS
s = base(); header(s, "ML Graphs & Analytics", "Results & Testing")
imgs = ["g1_learning_curve.png", "g3_accuracy_gauge.png",
        "g2_model_comparison.png", "g4_confidence.png"]
for i, fn in enumerate(imgs):
    col = i % 2; row = i // 2
    x = Inches(0.7 + col * 6.2); y = Inches(1.6 + row * 2.7)
    box(s, x - Inches(0.05), y - Inches(0.05), Inches(5.9), Inches(2.6),
        WHITE, MSO_SHAPE.RECTANGLE, line=RGBColor(0xE2,0xE8,0xF0), line_w=1)
    p = os.path.join(CLEAN, fn)
    if os.path.exists(p):
        s.shapes.add_picture(p, x, y, width=Inches(5.8))

# =========================================================== SLIDE 22 — TESTING
s = base(); header(s, "Testing", "Results & Testing")
tests = [
    ["ID", "Module", "Input", "Expected", "Result"],
    ["T01", "Auth", "Valid login", "JWT returned", "Pass"],
    ["T02", "Auth", "Wrong password", "401 error", "Pass"],
    ["T03", "OCR", "Prescription image", "Structured fields", "Pass"],
    ["T04", "ML Classify", "\"Amoxicillin\"", "Bacterial Infection", "Pass"],
    ["T05", "Interactions", "2 drugs", "Interaction list", "Pass"],
    ["T06", "Reminder", "Set 08:00", "Reminder saved", "Pass"],
    ["T07", "Health Score", "Vitals logged", "Score 0–100", "Pass"],
]
t = s.shapes.add_table(len(tests), 5, Inches(0.7), Inches(1.65),
                       Inches(11.9), Inches(5.1)).table
widths = [1.1, 2.2, 3.0, 3.4, 2.2]
for i, w in enumerate(widths):
    t.columns[i].width = Inches(w)
for ri, row in enumerate(tests):
    for ci, val in enumerate(row):
        cell = t.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci in (0, 4) else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val; r.font.name = BODYFONT; r.font.size = Pt(12.5)
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = SLATE
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else CARD
            if val == "Pass":
                r.font.color.rgb = GREEN; r.font.bold = True
            else:
                r.font.color.rgb = SLATE

# =========================================================== SLIDE 23 — FUTURE
s = base(); header(s, "Future Scope", "Conclusion")
ms = ["WebRTC live\nvideo consults", "Wearable sync\n(Apple/Google Fit)",
      "Predictive ML for\nvitals trends", "Expand to 50+\ndrug classes", "Mobile\napp"]
ln = s.shapes.add_connector(1, Inches(1.0), Inches(4.0), Inches(12.3), Inches(4.0))
ln.line.color.rgb = TEAL; ln.line.width = Pt(3); ln.shadow.inherit = False
for i, m in enumerate(ms):
    x = Inches(1.1 + i * 2.45)
    boxtext(s, x, Inches(3.7), Inches(0.6), Inches(0.6),
            BLUE if i % 2 else TEAL, [[(str(i+1), 16, WHITE, True, FONT)]],
            shape=MSO_SHAPE.OVAL)
    yy = Inches(2.3) if i % 2 == 0 else Inches(4.7)
    boxtext(s, x - Inches(0.7), yy, Inches(2.0), Inches(1.1), CARD,
            [[(m, 12.5, SLATE, True, BODYFONT)]])

# =========================================================== SLIDE 24 — THANK YOU
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, SW, SH, SLATE, MSO_SHAPE.RECTANGLE)
box(s, 0, 0, SW, Inches(0.25), TEAL, MSO_SHAPE.RECTANGLE)
box(s, 0, SH - Inches(0.25), SW, Inches(0.25), TEAL, MSO_SHAPE.RECTANGLE)
textbox(s, Inches(1), Inches(3.0), Inches(11.33), Inches(1.5),
        [[("Thank You", 60, HI, True, FONT)]], align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("Saved:", os.path.abspath(OUT))
print("Slides:", len(prs.slides._sldIdLst))
